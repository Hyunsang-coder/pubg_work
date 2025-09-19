from __future__ import annotations

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.group import GroupShape
from pptx.shapes.shapetree import SlideShapes
from pptx.shapes.autoshape import Shape
from pptx.shapes.picture import Picture
from pptx.chart.chart import Chart
from typing import List

from .models import SlideDoc, TextBlock, TableBlock, FigureBlock, NoteBlock


def _iter_shapes(shapes: SlideShapes):
    for shape in shapes:
        if isinstance(shape, GroupShape):
            for sub in _iter_shapes(shape.shapes):
                yield sub
        else:
            yield shape


def _shape_text_lines(shape: Shape) -> List[str]:
    lines: List[str] = []
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return lines
    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs)
        lines.append(text)
    return lines


def _shape_indent_levels(shape: Shape) -> List[int]:
    levels: List[int] = []
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return levels
    for paragraph in shape.text_frame.paragraphs:
        levels.append(paragraph.level or 0)
    return levels


def extract_slide(prs_slide, slide_index: int, options) -> SlideDoc:
    # title
    title = None
    if getattr(prs_slide.shapes, "title", None) is not None and prs_slide.shapes.title.has_text_frame:
        title = prs_slide.shapes.title.text.strip() or None
    if title is None:
        # fallback first textbox first line
        for shape in _iter_shapes(prs_slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and hasattr(shape, "text_frame") and shape.text_frame is not None:
                lines = _shape_text_lines(shape)
                if lines:
                    title = lines[0].strip()
                    break
    if title is None:
        title = f"Slide {slide_index+1}"

    slide_doc = SlideDoc(slide_index=slide_index, title=title, blocks=[])

    # collect blocks
    for shape in _iter_shapes(prs_slide.shapes):
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                rows = []
                tbl = shape.table
                for r in tbl.rows:
                    row_cells = []
                    for c in r.cells:
                        row_cells.append(c.text.strip())
                    rows.append(row_cells)
                slide_doc.blocks.append(TableBlock(shape_id=str(shape.shape_id), rows=rows, has_header=options.table_header))
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if options.figures != "omit":
                    slide_doc.blocks.append(FigureBlock(shape_id=str(shape.shape_id), figure_type="image", title=getattr(shape, "name", None)))
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                if options.charts == "labels":
                    # Chart text placeholders where available (pptx doesn't expose all easily)
                    chart: Chart = shape.chart  # type: ignore
                    title = None
                    if chart.has_title:
                        title = chart.chart_title.text_frame.text
                    slide_doc.blocks.append(FigureBlock(shape_id=str(shape.shape_id), figure_type="chart", title=title))
                elif options.figures == "placeholder":
                    slide_doc.blocks.append(FigureBlock(shape_id=str(shape.shape_id), figure_type="chart", title=None))
            else:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    lines = _shape_text_lines(shape)
                    if lines and any(line.strip() for line in lines):
                        slide_doc.blocks.append(
                            TextBlock(
                                shape_id=str(shape.shape_id),
                                lines=[l.strip() for l in lines if l is not None],
                                indent_levels=_shape_indent_levels(shape),
                            )
                        )
        except Exception:
            # skip problematic shapes but continue
            continue

    # notes
    if options.with_notes and hasattr(prs_slide, "notes_slide") and prs_slide.notes_slide and prs_slide.notes_slide.notes_text_frame:
        note_text = prs_slide.notes_slide.notes_text_frame.text
        if note_text.strip():
            slide_doc.blocks.append(NoteBlock(text=note_text.strip()))

    return slide_doc


def extract_pptx_to_docs(path: str, options) -> List[SlideDoc]:
    prs = Presentation(path)
    docs: List[SlideDoc] = []
    for idx, slide in enumerate(prs.slides):
        if options.slide_range and (idx + 1) not in options.slide_range:
            continue
        docs.append(extract_slide(slide, idx, options))
    return docs
