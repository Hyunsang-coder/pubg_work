from __future__ import annotations
from dataclasses import dataclass
from typing import List, Dict, Any, Callable, Optional
import io
import hashlib
import os
import zipfile
from pptx import Presentation
from pptx.shapes.group import GroupShape
from pptx.shapes.shapetree import SlideShapes
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL
from PIL import Image, ImageOps
from .translate import translate_texts, TranslationConfig


@dataclass
class ParagraphInfo:
    """단락 정보와 원본 서식을 저장하는 클래스"""
    slide_idx: int
    shape_id: str
    paragraph_idx: int
    original_text: str
    first_run_font: Dict[str, Any]  # 첫 번째 run의 서식 정보
    paragraph_ref: Any  # 실제 paragraph 객체 참조


def _iter_shapes(shapes: SlideShapes):
    """그룹 도형을 포함한 모든 도형을 재귀적으로 순회"""
    for s in shapes:
        if isinstance(s, GroupShape):
            for sub in _iter_shapes(s.shapes):
                yield sub
        else:
            yield s


def _extract_font_properties(run):
    """run의 폰트 속성을 딕셔너리로 추출"""
    font = run.font
    
    # 색상 처리 - 다양한 색상 타입에 대응
    color_info = None
    try:
        if font.color and hasattr(font.color, 'rgb') and font.color.rgb:
            color_info = {'type': 'rgb', 'value': font.color.rgb}
        elif font.color and hasattr(font.color, 'theme_color') and font.color.theme_color is not None:
            color_info = {'type': 'theme', 'value': font.color.theme_color}
    except Exception:
        color_info = None
    
    return {
        'name': font.name,
        'size': font.size,
        'bold': font.bold,
        'italic': font.italic,
        'underline': font.underline,
        'color': color_info,
    }


def _apply_font_properties(run, font_props: Dict[str, Any]):
    """저장된 폰트 속성을 run에 적용"""
    font = run.font
    
    # 기본 속성들 적용
    try:
        if font_props.get('name'):
            font.name = font_props['name']
    except Exception:
        pass
    
    try:
        if font_props.get('size'):
            font.size = font_props['size']
    except Exception:
        pass
    
    try:
        if font_props.get('bold') is not None:
            font.bold = font_props['bold']
    except Exception:
        pass
    
    try:
        if font_props.get('italic') is not None:
            font.italic = font_props['italic']
    except Exception:
        pass
    
    try:
        if font_props.get('underline') is not None:
            font.underline = font_props['underline']
    except Exception:
        pass
    
    # 색상 적용 - 안전하게 처리
    try:
        color_info = font_props.get('color')
        if color_info and isinstance(color_info, dict):
            if color_info['type'] == 'rgb' and color_info['value']:
                font.color.rgb = color_info['value']
            elif color_info['type'] == 'theme' and color_info['value'] is not None:
                font.color.theme_color = color_info['value']
    except Exception:
        # 색상 적용 실패 시 무시하고 계속 진행
        pass


def _iter_all_shapes(shapes: SlideShapes):
    """그룹 내 모든 도형 재귀 순회 (이미지 최적화용)"""
    for s in shapes:
        if isinstance(s, GroupShape):
            for sub in _iter_all_shapes(s.shapes):
                yield sub
        else:
            yield s


def _downscale_image(img: Image.Image, max_px: int) -> Image.Image:
    """긴 변 기준으로 리사이즈. max_px가 0이거나 작으면 그대로 반환"""
    if not max_px or max_px <= 0:
        return img
    width, height = img.size
    longest = max(width, height)
    if longest <= max_px:
        return img
    scale = max_px / float(longest)
    new_size = (max(1, int(width * scale)), max(1, int(height * scale)))
    return img.resize(new_size, Image.LANCZOS)


def _recompress_blob(original: bytes, *, quality: int, max_px: int) -> bytes | None:
    """이미지 바이트를 Pillow로 재압축. 용량이 줄어들지 않으면 None 반환."""
    try:
        Image.MAX_IMAGE_PIXELS = None
        with Image.open(io.BytesIO(original)) as im:
            im.load()
            # EXIF 방향 보정 (손상 방지)
            im = ImageOps.exif_transpose(im)
            fmt = (im.format or "PNG").upper()
            im2 = _downscale_image(im, max_px)

            buf = io.BytesIO()
            if fmt in ("JPEG", "JPG"):
                im3 = im2.convert("RGB")
                q = max(10, min(95, int(quality))) if isinstance(quality, int) else 70
                # baseline JPEG (progressive 비활성화)로 호환성 개선
                im3.save(buf, "JPEG", quality=q, optimize=True)
            elif fmt == "PNG":
                im2.save(buf, "PNG", optimize=True)
            else:
                return None
            data = buf.getvalue()
            if len(data) < len(original):
                return data
            return None
    except Exception:
        return None


def compress_images_in_presentation(
    prs: Presentation,
    *,
    quality: int = 70,
    max_px: int = 1920,
    progress_cb: Optional[Callable[[Dict[str, Any]], None]] = None,
) -> Dict[str, int]:
    """프레젠테이션 내 래스터 이미지(JPEG/PNG)를 다운스케일/재압축.

    Returns: {"pictures": 전체 그림 개수, "candidates": 처리 시도, "optimized": 최적화 성공, "failures": 실패, "bytes_saved": 절감 바이트}
    """
    cache: Dict[str, Optional[bytes]] = {}
    stats = {
        "pictures": 0,
        "fills": 0,
        "backgrounds": 0,
        "candidates": 0,
        "optimized": 0,
        "failures": 0,
        "bytes_saved": 0,
        "skipped_no_saving": 0,
        "unsupported_formats": 0,
    }
    if progress_cb:
        try:
            progress_cb({"message": "이미지 최적화 시작...", "ratio": 0.99})
        except Exception:
            pass
    for slide in prs.slides:
        for shp in _iter_all_shapes(slide.shapes):
            try:
                if shp.shape_type != MSO_SHAPE_TYPE.PICTURE or not hasattr(shp, "_element"):
                    # 도형이 이미지가 아니더라도 그림 채움(Picture Fill)인 경우 최적화 대상
                    try:
                        if hasattr(shp, "fill") and getattr(shp.fill, "type", None) == MSO_FILL.PICTURE:
                            stats["fills"] += 1
                            blipFill = getattr(shp.fill, "_xFill", None)
                            blipFill = getattr(blipFill, "blipFill", None)
                            if blipFill is None or getattr(blipFill, "blip", None) is None:
                                continue
                            blip = blipFill.blip
                            rId = getattr(blip, "embed", None) or getattr(blip, "link", None)
                            if not rId:
                                continue
                            part = shp.part.related_parts.get(rId)
                            if part is None or not hasattr(part, "blob"):
                                continue
                            orig = part.blob
                            stats["candidates"] += 1
                            h = hashlib.md5(orig).hexdigest()
                            if h in cache:
                                new_blob = cache[h]
                            else:
                                new_blob = _recompress_blob(orig, quality=quality, max_px=max_px)
                                cache[h] = new_blob
                            if new_blob and len(new_blob) < len(orig):
                                part._blob = new_blob  # type: ignore[attr-defined]
                                stats["optimized"] += 1
                                stats["bytes_saved"] += (len(orig) - len(new_blob))
                                if progress_cb and stats["optimized"] % 10 == 0:
                                    try:
                                        progress_cb({"message": f"이미지 최적화 진행 중... ({stats['optimized']}개)", "ratio": 0.995})
                                    except Exception:
                                        pass
                            else:
                                stats["skipped_no_saving"] += 1
                        continue
                    except Exception:
                        stats["failures"] += 1
                        continue
                stats["pictures"] += 1
                blipFill = getattr(shp._element, "blipFill", None)
                if blipFill is None or getattr(blipFill, "blip", None) is None:
                    continue
                blip = blipFill.blip
                rId = getattr(blip, "embed", None) or getattr(blip, "link", None)
                if not rId:
                    continue
                part = shp.part.related_parts.get(rId)
                if part is None or not hasattr(part, "blob"):
                    continue
                orig = part.blob
                stats["candidates"] += 1
                h = hashlib.md5(orig).hexdigest()
                if h in cache:
                    new_blob = cache[h]
                else:
                    new_blob = _recompress_blob(orig, quality=quality, max_px=max_px)
                    cache[h] = new_blob
                if new_blob and len(new_blob) < len(orig):
                    part._blob = new_blob  # type: ignore[attr-defined]
                    stats["optimized"] += 1
                    stats["bytes_saved"] += (len(orig) - len(new_blob))
                    if progress_cb and stats["optimized"] % 10 == 0:
                        try:
                            progress_cb({"message": f"이미지 최적화 진행 중... ({stats['optimized']}개)", "ratio": 0.995})
                        except Exception:
                            pass
                else:
                    stats["skipped_no_saving"] += 1
            except Exception:
                stats["failures"] += 1
                continue
        # 슬라이드 배경 그림 최적화
        try:
            if hasattr(slide, "background") and hasattr(slide.background, "fill"):
                fill = slide.background.fill
                if getattr(fill, "type", None) == MSO_FILL.PICTURE:
                    stats["backgrounds"] += 1
                    blipFill = getattr(fill, "_xFill", None)
                    blipFill = getattr(blipFill, "blipFill", None)
                    if blipFill is not None and getattr(blipFill, "blip", None) is not None:
                        blip = blipFill.blip
                        rId = getattr(blip, "embed", None) or getattr(blip, "link", None)
                        if rId:
                            part = slide.part.related_parts.get(rId)
                            if part is not None and hasattr(part, "blob"):
                                orig = part.blob
                                stats["candidates"] += 1
                                h = hashlib.md5(orig).hexdigest()
                                if h in cache:
                                    new_blob = cache[h]
                                else:
                                    new_blob = _recompress_blob(orig, quality=quality, max_px=max_px)
                                    cache[h] = new_blob
                                if new_blob and len(new_blob) < len(orig):
                                    part._blob = new_blob  # type: ignore[attr-defined]
                                    stats["optimized"] += 1
                                    stats["bytes_saved"] += (len(orig) - len(new_blob))
                                else:
                                    stats["skipped_no_saving"] += 1
        except Exception:
            stats["failures"] += 1
    if progress_cb:
        try:
            mb = stats["bytes_saved"] / (1024 * 1024)
            progress_cb({"message": f"이미지 최적화 완료 — 후보 {stats['candidates']}개, 성공 {stats['optimized']}개, 절감 {mb:.1f}MB (채움 {stats['fills']} / 배경 {stats['backgrounds']})", "ratio": 0.999})
        except Exception:
            pass
    return stats


def optimize_pptx_media_zip(
    input_pptx_path: str,
    output_pptx_path: str,
    *,
    quality: int = 70,
    max_px: int = 1920,
    progress_cb: Optional[Callable[[Dict[str, Any]], None]] = None,
) -> Dict[str, int]:
    """ZIP 레벨에서 ppt/media/* 이미지를 재압축하여 XML 손상 리스크 없이 용량 최적화.

    Returns: {"media": 총 media 이미지 수, "optimized": 성공, "failures": 실패, "bytes_saved": 절감 바이트}
    """
    stats = {"media": 0, "optimized": 0, "failures": 0, "bytes_saved": 0}
    if progress_cb:
        try:
            progress_cb({"message": "미디어 스캔 시작...", "ratio": 0.01})
        except Exception:
            pass
    os.makedirs(os.path.dirname(os.path.abspath(output_pptx_path)), exist_ok=True)
    with zipfile.ZipFile(input_pptx_path, "r") as zin, zipfile.ZipFile(output_pptx_path, "w", compression=zipfile.ZIP_DEFLATED) as zout:
        namelist = zin.namelist()
        total_media = sum(1 for n in namelist if n.startswith("ppt/media/") and n.lower().endswith((".jpg", ".jpeg", ".png")))
        processed = 0
        for name in namelist:
            data = zin.read(name)
            lower = name.lower()
            if name.startswith("ppt/media/") and lower.endswith((".jpg", ".jpeg", ".png")):
                stats["media"] += 1
                try:
                    new_blob = _recompress_blob(data, quality=quality, max_px=max_px)
                    if new_blob and len(new_blob) < len(data):
                        zout.writestr(name, new_blob)
                        stats["optimized"] += 1
                        stats["bytes_saved"] += (len(data) - len(new_blob))
                    else:
                        zout.writestr(name, data)
                    processed += 1
                    if progress_cb and total_media:
                        ratio = 0.02 + 0.96 * (processed / total_media)
                        try:
                            progress_cb({"message": f"미디어 최적화 진행 {processed}/{total_media}", "ratio": ratio})
                        except Exception:
                            pass
                except Exception:
                    stats["failures"] += 1
                    zout.writestr(name, data)
            else:
                zout.writestr(name, data)
    return stats


def create_translated_presentation_v2(
    input_pptx: str,
    output_pptx: str,
    config: TranslationConfig,
    *,
    progress_callback: Optional[Callable[[str], None]] = None,
    batch_size: int = 400,
    image_opt: Optional[Dict[str, Any]] = None,
) -> Dict[str, int]:
    """
    하이브리드 접근 방식으로 프레젠테이션을 번역하는 함수
    
    Args:
        input_pptx: 입력 PPTX 파일 경로
        output_pptx: 출력 PPTX 파일 경로
        config: 번역 설정

    Returns:
        dict: 번역된 프레젠테이션의 통계 정보 (슬라이드 수, 번역 단어 수)
    """
    # 1. 프레젠테이션 로드
    prs = Presentation(input_pptx)
    slide_count = len(prs.slides)
    
    # 2. 모든 단락 정보 수집
    # 진행 로그: 텍스트 및 서식 수집 시작
    if progress_callback:
        try:
            progress_callback({"message": "PPT 분석: 텍스트 및 서식 수집 중...", "ratio": 0.02})
        except Exception:
            pass
    paragraph_infos: List[ParagraphInfo] = []
    texts_to_translate: List[str] = []
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape in _iter_shapes(slide.shapes):
            # 표 도형은 셀 단위로 텍스트를 별도로 다룬다
            if getattr(shape, 'has_table', False) and shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if not cell.text_frame:
                            continue
                        for para_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                            if not paragraph.runs:
                                continue

                            full_paragraph_text = "".join(run.text for run in paragraph.runs)
                            if not full_paragraph_text.strip():
                                continue

                            first_run_font = _extract_font_properties(paragraph.runs[0])
                            para_info = ParagraphInfo(
                                slide_idx=slide_idx,
                                shape_id=str(shape.shape_id),
                                paragraph_idx=para_idx,
                                original_text=full_paragraph_text,
                                first_run_font=first_run_font,
                                paragraph_ref=paragraph,
                            )

                            paragraph_infos.append(para_info)
                            texts_to_translate.append(full_paragraph_text)
                continue

            if not hasattr(shape, 'text_frame') or not shape.text_frame:
                continue

            text_frame = shape.text_frame
            for para_idx, paragraph in enumerate(text_frame.paragraphs):
                if not paragraph.runs:
                    continue

                # 3. 단락 단위 텍스트 통합
                full_paragraph_text = "".join(run.text for run in paragraph.runs)

                if not full_paragraph_text.strip():
                    continue

                # 첫 번째 run의 서식 정보 저장
                first_run_font = _extract_font_properties(paragraph.runs[0])

                # 단락 정보 저장
                para_info = ParagraphInfo(
                    slide_idx=slide_idx,
                    shape_id=str(shape.shape_id),
                    paragraph_idx=para_idx,
                    original_text=full_paragraph_text,
                    first_run_font=first_run_font,
                    paragraph_ref=paragraph
                )

                paragraph_infos.append(para_info)
                texts_to_translate.append(full_paragraph_text)
    
    def _log(message: str, *, ratio: float | None = None) -> None:
        if progress_callback:
            payload = {"message": message}
            if ratio is not None:
                payload["ratio"] = max(0.0, min(1.0, ratio))
            try:
                progress_callback(payload)
            except Exception:
                pass

    # 4. 일괄 번역
    word_count_source = sum(len(text.split()) for text in texts_to_translate)

    if not texts_to_translate:
        # 번역할 텍스트가 없으면 원본 복사
        _log("번역할 문장이 없어 원본 PPT를 그대로 저장합니다.", ratio=1.0)
        prs.save(output_pptx)
        return {"slides": slide_count, "word_count": 0}

    total = len(texts_to_translate)
    _log(f"슬라이드 분석 완료: 총 {total}개 문장", ratio=0.05)
    translated_texts: List[str] = []
    safe_batch_size = max(1, batch_size)
    progress_base = 0.1
    progress_span = 0.7

    total_batches = (total + safe_batch_size - 1) // safe_batch_size if total else 0

    for start in range(0, total, safe_batch_size):
        end = min(start + safe_batch_size, total)
        batch = texts_to_translate[start:end]
        completed = end
        ratio = progress_base + progress_span * (completed / total)
        current_batch = (start // safe_batch_size) + 1
        batch_caption = f"배치 {current_batch}/{total_batches}" if total_batches else "배치 진행"
        pre_ratio = max(progress_base, ratio - 0.02)
        _log(f"번역 요청 준비 — {batch_caption} ({start + 1}~{end}/{total})", ratio=pre_ratio)
        translated_batch = translate_texts(batch, config)
        _log(f"모델 응답 수신 — {batch_caption} ({start + 1}~{end}/{total})", ratio=ratio)
        translated_texts.extend(translated_batch)

    _log("PPT 반영: 번역된 텍스트 재삽입 및 서식 복원 중...", ratio=0.9)
    
    # 5. 번역된 텍스트 재삽입
    total_apply = len(paragraph_infos)
    for idx, (para_info, translated_text) in enumerate(zip(paragraph_infos, translated_texts)):
        if total_apply > 0 and (idx + 1) % 50 == 0:
            # 0.90~0.99 구간에서 세밀한 진행률 업데이트
            stage_progress = (idx + 1) / max(1, total_apply)
            _log(f"PPT 반영 진행 중... ({idx + 1}/{total_apply})", ratio=0.9 + 0.09 * stage_progress)
        paragraph = para_info.paragraph_ref
        
        # 안전성 검사
        if not translated_text or not isinstance(translated_text, str):
            translated_text = para_info.original_text
        
        try:
            # 더 안전한 방법: 모든 run을 삭제하고 새로 생성
            # 기존 서식 정보 백업
            original_font_props = para_info.first_run_font
            
            # paragraph의 모든 텍스트를 새 텍스트로 교체
            paragraph.clear()
            new_run = paragraph.add_run()
            new_run.text = translated_text
            
            # 원본 서식 복원
            _apply_font_properties(new_run, original_font_props)
            
        except Exception:
            # 개별 단락 처리 실패 시 무시하고 계속 진행
            continue
    translated_word_count = sum(len(text.split()) for text in translated_texts) if translated_texts else word_count_source

    # 선택적 이미지 최적화 단계
    if image_opt:
        try:
            q = int(image_opt.get("quality", 70))  # type: ignore[union-attr]
            max_px = int(image_opt.get("max_px", 1920))  # type: ignore[union-attr]
            _log("이미지 최적화 중...", ratio=0.99)
            img_stats = compress_images_in_presentation(prs, quality=q, max_px=max_px, progress_cb=progress_callback)
            try:
                saved_mb = float(img_stats.get("bytes_saved", 0)) / (1024 * 1024)
                _log(f"이미지 최적화 요약 — 성공 {img_stats.get('optimized',0)}개, 절감 {saved_mb:.1f}MB", ratio=0.999)
            except Exception:
                pass
        except Exception:
            # 최적화 실패해도 전체 흐름에는 영향 주지 않음
            pass

    _log("번역 결과 저장 완료", ratio=1.0)
    # 6. 프레젠테이션 저장
    prs.save(output_pptx)
    return {"slides": slide_count, "word_count": translated_word_count}


# 레거시 호환성을 위한 래퍼 (사용하지 않음)
# def create_translated_copy_v2(...): 삭제됨
