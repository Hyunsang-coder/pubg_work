from __future__ import annotations
from dataclasses import dataclass
from typing import List, Tuple, Dict, Any
from pptx import Presentation
from pptx.shapes.group import GroupShape
from pptx.shapes.shapetree import SlideShapes
from .translate import translate_texts, TranslationConfig


@dataclass
class ParagraphInfo:
    """단락 정보와 원본 서식을 저장하는 클래스"""
    slide_idx: int
    shape_id: str
    paragraph_idx: int
    original_text: str
    first_run_font: Dict[str, Any]  # 첫 번째 run의 서식 정보
    paragraph_ref: Any  # 실제 paragraph 객체 참조


def _iter_shapes(shapes: SlideShapes):
    """그룹 도형을 포함한 모든 도형을 재귀적으로 순회"""
    for s in shapes:
        if isinstance(s, GroupShape):
            for sub in _iter_shapes(s.shapes):
                yield sub
        else:
            yield s


def _extract_font_properties(run):
    """run의 폰트 속성을 딕셔너리로 추출"""
    font = run.font
    
    # 색상 처리 - 다양한 색상 타입에 대응
    color_info = None
    try:
        if font.color and hasattr(font.color, 'rgb') and font.color.rgb:
            color_info = {'type': 'rgb', 'value': font.color.rgb}
        elif font.color and hasattr(font.color, 'theme_color') and font.color.theme_color is not None:
            color_info = {'type': 'theme', 'value': font.color.theme_color}
    except Exception:
        color_info = None
    
    return {
        'name': font.name,
        'size': font.size,
        'bold': font.bold,
        'italic': font.italic,
        'underline': font.underline,
        'color': color_info,
    }


def _apply_font_properties(run, font_props: Dict[str, Any]):
    """저장된 폰트 속성을 run에 적용"""
    font = run.font
    
    # 기본 속성들 적용
    try:
        if font_props.get('name'):
            font.name = font_props['name']
    except Exception:
        pass
    
    try:
        if font_props.get('size'):
            font.size = font_props['size']
    except Exception:
        pass
    
    try:
        if font_props.get('bold') is not None:
            font.bold = font_props['bold']
    except Exception:
        pass
    
    try:
        if font_props.get('italic') is not None:
            font.italic = font_props['italic']
    except Exception:
        pass
    
    try:
        if font_props.get('underline') is not None:
            font.underline = font_props['underline']
    except Exception:
        pass
    
    # 색상 적용 - 안전하게 처리
    try:
        color_info = font_props.get('color')
        if color_info and isinstance(color_info, dict):
            if color_info['type'] == 'rgb' and color_info['value']:
                font.color.rgb = color_info['value']
            elif color_info['type'] == 'theme' and color_info['value'] is not None:
                font.color.theme_color = color_info['value']
    except Exception:
        # 색상 적용 실패 시 무시하고 계속 진행
        pass


def create_translated_presentation_v2(
    input_pptx: str, 
    output_pptx: str, 
    config: TranslationConfig
) -> None:
    """
    하이브리드 접근 방식으로 프레젠테이션을 번역하는 함수
    
    Args:
        input_pptx: 입력 PPTX 파일 경로
        output_pptx: 출력 PPTX 파일 경로
        config: 번역 설정
    """
    # 1. 프레젠테이션 로드
    prs = Presentation(input_pptx)
    
    # 2. 모든 단락 정보 수집
    paragraph_infos: List[ParagraphInfo] = []
    texts_to_translate: List[str] = []
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape in _iter_shapes(slide.shapes):
            if not hasattr(shape, 'text_frame') or not shape.text_frame:
                continue
                
            text_frame = shape.text_frame
            for para_idx, paragraph in enumerate(text_frame.paragraphs):
                if not paragraph.runs:
                    continue
                    
                # 3. 단락 단위 텍스트 통합
                full_paragraph_text = "".join(run.text for run in paragraph.runs)
                
                if not full_paragraph_text.strip():
                    continue
                
                # 디버깅: 추출된 텍스트 로깅
                print(f"📝 슬라이드 {slide_idx}, Shape {shape.shape_id}: '{full_paragraph_text[:50]}{'...' if len(full_paragraph_text) > 50 else ''}'")
                
                # 첫 번째 run의 서식 정보 저장
                first_run_font = _extract_font_properties(paragraph.runs[0])
                
                # 단락 정보 저장
                para_info = ParagraphInfo(
                    slide_idx=slide_idx,
                    shape_id=str(shape.shape_id),
                    paragraph_idx=para_idx,
                    original_text=full_paragraph_text,
                    first_run_font=first_run_font,
                    paragraph_ref=paragraph
                )
                
                paragraph_infos.append(para_info)
                texts_to_translate.append(full_paragraph_text)
    
    # 4. 일괄 번역
    if not texts_to_translate:
        # 번역할 텍스트가 없으면 원본 복사
        prs.save(output_pptx)
        return
    
    print(f"📊 번역 통계:")
    print(f"  - 총 슬라이드 수: {len(prs.slides)}")
    print(f"  - 추출된 단락 수: {len(paragraph_infos)}")
    print(f"  - 번역할 텍스트 수: {len(texts_to_translate)}")
    print(f"  - 샘플 텍스트: {texts_to_translate[:3] if texts_to_translate else '없음'}")
    
    translated_texts = translate_texts(texts_to_translate, config)
    
    print(f"  - 번역된 텍스트 수: {len(translated_texts)}")
    print(f"  - 샘플 번역: {translated_texts[:3] if translated_texts else '없음'}")
    
    # 5. 번역된 텍스트 재삽입
    for para_info, translated_text in zip(paragraph_infos, translated_texts):
        paragraph = para_info.paragraph_ref
        
        # 안전성 검사
        if not translated_text or not isinstance(translated_text, str):
            translated_text = para_info.original_text
        
        try:
            # 더 안전한 방법: 모든 run을 삭제하고 새로 생성
            # 기존 서식 정보 백업
            original_font_props = para_info.first_run_font
            
            # paragraph의 모든 텍스트를 새 텍스트로 교체
            paragraph.clear()
            new_run = paragraph.add_run()
            new_run.text = translated_text
            
            # 원본 서식 복원
            _apply_font_properties(new_run, original_font_props)
            
        except Exception as e:
            # 개별 단락 처리 실패 시 로그 남기고 계속 진행
            print(f"⚠️ Warning: Failed to process paragraph in slide {para_info.slide_idx}, shape {para_info.shape_id}: {e}")
            continue
    
    # 6. 프레젠테이션 저장
    prs.save(output_pptx)


# 기존 함수와의 호환성을 위한 래퍼
def create_translated_copy_v2(
    input_pptx: str, 
    translated_docs, 
    output_pptx: str, 
    policy, 
    config: TranslationConfig
):
    """기존 인터페이스와 호환되는 래퍼 함수"""
    create_translated_presentation_v2(input_pptx, output_pptx, config)