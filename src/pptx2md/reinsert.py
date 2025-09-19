from __future__ import annotations
from dataclasses import dataclass
from typing import List
from pptx import Presentation
from pptx.util import Pt
from pptx.shapes.group import GroupShape
from pptx.shapes.shapetree import SlideShapes
from .models import SlideDoc, TextBlock, TableBlock
from .translate import shorten_line

@dataclass
class OverflowPolicy:
    max_chars_per_paragraph: int = 180
    min_font_size_pt: int = 12
    reduce_step_pt: int = 1

def _reduce_font_size(tf, min_pt: int, step_pt: int):
    for p in tf.paragraphs:
        for run in p.runs:
            size = run.font.size.pt if run.font.size and run.font.size.pt else 18
            new_size = max(min_pt, int(size) - step_pt)
            run.font.size = Pt(new_size)

def _apply_text_to_shape(shape, lines: List[str], policy: OverflowPolicy):
    tf = shape.text_frame
    tf.clear()
    # sanitize and shorten
    safe_lines: List[str] = []
    for line in (lines or []):
        text = "" if line is None else str(line)
        text = shorten_line(text, policy.max_chars_per_paragraph)
        safe_lines.append(text)
    if not safe_lines:
        safe_lines = [""]
    for i, text in enumerate(safe_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        # preserve bullet indent if available: if original had indent levels we'd need them; absent here, keep default
        # could be improved by receiving indent levels in block
    if any(len(l or "") > policy.max_chars_per_paragraph for l in safe_lines):
        _reduce_font_size(tf, policy.min_font_size_pt, policy.reduce_step_pt)

def _apply_text_to_table(shape, rows: List[List[str]], policy: OverflowPolicy):
    tbl = shape.table
    for r, row_data in enumerate(rows or []):
        for c, cell_text in enumerate(row_data or []):
            if r < len(tbl.rows) and c < len(tbl.columns):
                safe = "" if cell_text is None else str(cell_text)
                tbl.cell(r, c).text = shorten_line(safe, policy.max_chars_per_paragraph)

def _iter_shapes(shapes: SlideShapes):
    for s in shapes:
        if isinstance(s, GroupShape):
            for sub in _iter_shapes(s.shapes):
                yield sub
        else:
            yield s


def create_translated_copy(input_pptx: str, translated_docs: List[SlideDoc], output_pptx: str, policy: OverflowPolicy):
    prs = Presentation(input_pptx)
    translation_map = {
        doc.slide_index: {
            block.shape_id: block for block in doc.blocks if hasattr(block, "shape_id")
        } for doc in translated_docs
    }
    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx not in translation_map:
            continue
        slide_map = translation_map[slide_idx]
        for shape in _iter_shapes(slide.shapes):
            shape_id = str(shape.shape_id)
            if shape_id in slide_map:
                block = slide_map[shape_id]
                if isinstance(block, TextBlock) and hasattr(shape, "text_frame"):
                    _apply_text_to_shape(shape, block.lines, policy)
                elif isinstance(block, TableBlock) and hasattr(shape, "table"):
                    _apply_text_to_table(shape, block.rows, policy)
    prs.save(output_pptx)
